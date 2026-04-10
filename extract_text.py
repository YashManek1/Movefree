import os
from PyPDF2 import PdfReader
from pptx import Presentation

os.makedirs('scratch', exist_ok=True)

def extract_pdf(pdf_path, txt_path):
    try:
        reader = PdfReader(pdf_path)
        text = ''
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + '\n'
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"Extracted {pdf_path} to {txt_path}")
    except Exception as e:
        print(f"Error extracting {pdf_path}: {e}")

def extract_pptx(pptx_path, txt_path):
    try:
        prs = Presentation(pptx_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"Extracted {pptx_path} to {txt_path}")
    except Exception as e:
        print(f"Error extracting {pptx_path}: {e}")

extract_pdf('MoveFree_Research_Paper_final.pdf', 'scratch/paper_text.txt')
extract_pdf('IPD_Sem5_Report.pdf', 'scratch/report_text.txt')
extract_pptx('IPD_Sem_5.pptx', 'scratch/ppt_text.txt')
